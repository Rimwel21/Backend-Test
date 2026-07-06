from pathlib import Path
import re
import shutil
import zipfile
from uuid import uuid4
from fastapi.responses import FileResponse
from fastapi import HTTPException, Request, UploadFile, status
from sqlalchemy.orm import Session, joinedload
from models.accounts import Accounts
from models.learning_topic import LearningTopic
from models.teacher_class import TeacherClass
from models.teacher_module import TeacherModule
from schemas.teacher_module_schema import TeacherModuleCreate, TeacherModuleUpdate
from utils.enum import RoleEnum

ALLOWED_MATERIAL_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".doc", ".docx"}
ALLOWED_MATERIAL_TYPES = {
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
UPLOAD_DIR = Path("uploads/learning_materials")
MATERIAL_IMAGE_DIR = Path("static/material_images")
PDF_PAGE_DIR = Path("static/pdf_pages")


def _ensure_teacher(current_user: Accounts):
    if current_user.role != RoleEnum.teacher:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Teacher only")


def _validate_class_owner(class_id: int | None, db: Session, current_user: Accounts):
    if class_id is None:
        return

    existing_class = db.query(TeacherClass).filter(
        TeacherClass.id == class_id,
        TeacherClass.teacher_id == current_user.id
    ).first()

    if not existing_class:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Class not found")


def _validate_class_assignment(class_id: int | None, status_value: str):
    if status_value == "Published" and class_id is None:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Published materials must be assigned to a class")


def list_teacher_modules(request: Request, db: Session, current_user: Accounts):
    _ensure_teacher(current_user)

    return (
        db.query(TeacherModule)
        .filter(TeacherModule.teacher_id == current_user.id)
        .order_by(TeacherModule.created_at.desc())
        .all()
    )


def create_teacher_module(request: Request, module: TeacherModuleCreate, db: Session, current_user: Accounts):
    _ensure_teacher(current_user)
    _validate_class_assignment(module.class_id, module.status)
    _validate_class_owner(module.class_id, db, current_user)

    new_module = TeacherModule(
        teacher_id=current_user.id,
        class_id=module.class_id,
        title=module.title.strip(),
        description=module.description.strip(),
        content_type=module.content_type.strip() if module.content_type else None,
        week=module.week.strip() if module.week else None,
        file_name=module.file_name.strip() if module.file_name else None,
        file_type=module.file_type.strip() if module.file_type else None,
        status=module.status,
        behavior_required=str(module.behavior_required).lower(),
        estimated_time=module.estimated_time.strip() if module.estimated_time else None,
    )

    db.add(new_module)
    db.commit()
    db.refresh(new_module)
    _replace_generated_topics(db, new_module, _build_fallback_text(new_module.title, new_module.description, new_module.file_name), [])
    db.refresh(new_module)

    return new_module


async def create_teacher_module_upload(
    request: Request,
    title: str,
    description: str,
    content_type: str,
    week: str,
    status_value: str,
    behavior_required: bool,
    estimated_time: str | None,
    class_id: int | None,
    material_file: UploadFile,
    db: Session,
    current_user: Accounts
):
    _ensure_teacher(current_user)
    _validate_class_assignment(class_id, status_value)
    _validate_class_owner(class_id, db, current_user)
    _validate_material_file(material_file)
    saved_path, file_size = _save_material_file(material_file)

    new_module = TeacherModule(
        teacher_id=current_user.id,
        class_id=class_id,
        title=title.strip(),
        description=description.strip(),
        content_type=content_type.strip(),
        week=week.strip(),
        file_name=material_file.filename,
        file_type=material_file.content_type,
        file_path=str(saved_path),
        file_size=file_size,
        status=status_value,
        behavior_required=str(behavior_required).lower(),
        estimated_time=estimated_time.strip() if estimated_time else None,
    )

    db.add(new_module)
    db.commit()
    db.refresh(new_module)
    if saved_path.suffix.lower() == ".pdf":
        _replace_topic_records(db, new_module, _extract_pdf_logical_topics(saved_path, new_module.title, new_module.description))
    else:
        extracted_text = _extract_material_text(saved_path, new_module.title, new_module.description)
        image_urls = _extract_material_images(saved_path)
        _replace_generated_topics(db, new_module, extracted_text, image_urls)
    db.refresh(new_module)

    return new_module


def get_teacher_module(request: Request, module_id: int, db: Session, current_user: Accounts):
    _ensure_teacher(current_user)

    module = db.query(TeacherModule).options(joinedload(TeacherModule.topics)).filter(
        TeacherModule.id == module_id,
        TeacherModule.teacher_id == current_user.id
    ).first()

    if not module:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Module not found")

    if not module.topics:
        _replace_generated_topics(db, module, _build_fallback_text(module.title, module.description, module.file_name), [])
        db.refresh(module)
    elif _needs_pdf_topic_regeneration(module):
        _replace_topic_records(db, module, _extract_pdf_logical_topics(Path(module.file_path), module.title, module.description))
        db.refresh(module)

    return module


def update_teacher_module(request: Request, module_id: int, update: TeacherModuleUpdate, db: Session, current_user: Accounts):
    module = get_teacher_module(request, module_id, db, current_user)
    update_data = update.model_dump(exclude_unset=True)
    next_class_id = update_data.get("class_id", module.class_id)
    next_status = update_data.get("status", module.status)
    _validate_class_assignment(next_class_id, next_status)

    if "class_id" in update_data:
        _validate_class_owner(update.class_id, db, current_user)

    for key, value in update_data.items():
        if key == "behavior_required" and isinstance(value, bool):
            setattr(module, key, str(value).lower())
        else:
            setattr(module, key, value.strip() if isinstance(value, str) else value)

    db.commit()
    db.refresh(module)
    if any(key in update_data for key in {"title", "description"}):
        if not module.topics:
            _replace_generated_topics(db, module, _build_fallback_text(module.title, module.description, module.file_name), [])
            db.refresh(module)

    return module


def delete_teacher_module(request: Request, module_id: int, db: Session, current_user: Accounts):
    module = get_teacher_module(request, module_id, db, current_user)
    file_path = Path(module.file_path) if module.file_path else None

    db.delete(module)
    db.commit()
    if file_path and file_path.exists():
        file_path.unlink()

    return {"detail": "Module deleted successfully"}


async def replace_teacher_module_file(request: Request, module_id: int, material_file: UploadFile, db: Session, current_user: Accounts):
    module = get_teacher_module(request, module_id, db, current_user)
    _validate_material_file(material_file)

    old_path = Path(module.file_path) if module.file_path else None
    saved_path, file_size = _save_material_file(material_file)
    module.file_name = material_file.filename
    module.file_type = material_file.content_type
    module.file_path = str(saved_path)
    module.file_size = file_size

    if saved_path.suffix.lower() == ".pdf":
        _replace_topic_records(db, module, _extract_pdf_logical_topics(saved_path, module.title, module.description))
    else:
        _replace_generated_topics(db, module, _extract_material_text(saved_path, module.title, module.description), _extract_material_images(saved_path))
    db.commit()
    db.refresh(module)

    if old_path and old_path.exists() and old_path != saved_path:
        old_path.unlink()

    return module


def download_teacher_module_file(request: Request, module_id: int, db: Session, current_user: Accounts):
    module = get_teacher_module(request, module_id, db, current_user)
    if not module.file_path:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="No uploaded file is attached to this material")

    path = Path(module.file_path)
    if not path.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Uploaded file was not found on the server")

    return FileResponse(path, media_type=module.file_type or "application/octet-stream", filename=module.file_name or path.name)


def _validate_material_file(material_file: UploadFile):
    filename = material_file.filename or ""
    extension = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if extension not in ALLOWED_MATERIAL_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Unsupported file type. Upload PDF, PowerPoint, or Word files only."
        )

    if material_file.content_type and material_file.content_type not in ALLOWED_MATERIAL_TYPES:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Unsupported file type. Upload PDF, PowerPoint, or Word files only."
        )


def _save_material_file(material_file: UploadFile):
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    original_name = Path(material_file.filename or "material").name
    safe_name = f"{uuid4().hex}_{original_name}"
    destination = UPLOAD_DIR / safe_name
    material_file.file.seek(0)

    with destination.open("wb") as buffer:
        shutil.copyfileobj(material_file.file, buffer)

    material_file.file.seek(0)
    return destination, destination.stat().st_size


def _replace_generated_topics(db: Session, module: TeacherModule, extracted_text: str, image_urls: list[str] | None = None):
    _replace_topic_records(db, module, _split_topics(module.title, module.description, extracted_text, image_urls or []))


def _replace_topic_records(db: Session, module: TeacherModule, topics: list[dict]):
    module.topics.clear()
    db.flush()

    for index, topic in enumerate(topics, start=1):
        module.topics.append(
            LearningTopic(
                title=topic["title"],
                description=topic["description"],
                content=topic["content"],
                image_url=topic.get("image_url"),
                page_image_urls=topic.get("page_image_urls", []),
                sort_order=index,
            )
        )

    db.commit()


def _needs_pdf_topic_regeneration(module: TeacherModule):
    if not module.file_path or Path(module.file_path).suffix.lower() != ".pdf":
        return False
    if not Path(module.file_path).exists():
        return False
    if not module.topics:
        return True
    return not any(topic.page_image_urls for topic in module.topics)


def _split_topics(title: str, description: str, extracted_text: str, image_urls: list[str] | None = None):
    paragraphs = [line.strip() for line in extracted_text.splitlines() if line.strip()]
    if not paragraphs:
        paragraphs = [description.strip() or f"Open the uploaded file for {title}."]

    images = image_urls or []
    topics = [{
        "title": "Introduction",
        "description": description.strip() or title.strip(),
        "content": "\n\n".join(paragraphs[:4]) or description.strip(),
        "image_url": images[0] if images else None,
        "page_image_urls": [],
    }]

    remaining = paragraphs[4:]
    chunk_size = 5
    for start in range(0, len(remaining), chunk_size):
        chunk = remaining[start:start + chunk_size]
        heading = _topic_heading(chunk[0], len(topics))
        topics.append({
            "title": heading,
            "description": chunk[0][:180],
            "content": "\n\n".join(chunk),
            "image_url": images[len(topics)] if len(topics) < len(images) else None,
            "page_image_urls": [],
        })

    if len(topics) == 1 and len(paragraphs) > 1:
        topics[0]["content"] = "\n\n".join(paragraphs)

    return topics


def _topic_heading(candidate: str, topic_index: int):
    clean = candidate.strip().strip(":")
    if 4 <= len(clean) <= 70 and len(clean.split()) <= 10:
        return clean
    return f"Topic {topic_index}"


def _build_fallback_text(title: str, description: str, file_name: str | None):
    file_line = f"Uploaded file: {file_name}" if file_name else "No file has been uploaded yet."
    return f"{title}\n\n{description}\n\n{file_line}"


def _extract_pdf_logical_topics(path: Path, module_title: str, module_description: str):
    try:
        import fitz
    except ImportError:
        return _split_topics(module_title, module_description, _extract_pdf_text(path), [])

    try:
        document = fitz.open(str(path))
    except Exception:
        return _split_topics(module_title, module_description, _extract_material_text(path, module_title, module_description), [])
    if document.page_count == 0:
        return _split_topics(module_title, module_description, _build_fallback_text(module_title, module_description, path.name), [])

    page_infos = [_pdf_page_info(document, page_index) for page_index in range(document.page_count)]
    boundaries = _detect_pdf_topic_boundaries(page_infos)
    if not boundaries:
        boundaries = _fallback_pdf_boundaries(page_infos)

    topics = []
    for topic_number, boundary in enumerate(boundaries, start=1):
        start = boundary["page_index"]
        end = boundaries[topic_number]["page_index"] if topic_number < len(boundaries) else document.page_count
        pages = page_infos[start:end]
        content = "\n\n".join(page["text"] for page in pages if page["text"]).strip()
        if not content:
            content = f"Open the uploaded PDF page group for {boundary['title']}."
        topics.append({
            "title": boundary["title"],
            "description": _summarize_topic(boundary["title"], content),
            "content": content,
            "image_url": None,
            "page_image_urls": _render_pdf_topic_pages(document, start, end),
        })

    document.close()
    return topics or _split_topics(module_title, module_description, _extract_pdf_text(path), [])


def _pdf_page_info(document, page_index: int):
    page = document.load_page(page_index)
    text = page.get_text("text", sort=True).strip()
    lines = []
    raw = page.get_text("dict", sort=True)
    sizes = []

    for block in raw.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = [span for span in line.get("spans", []) if span.get("text", "").strip()]
            if not spans:
                continue
            line_text = _normalize_pdf_line(" ".join(span["text"] for span in spans))
            if not line_text:
                continue
            size = max(float(span.get("size", 0)) for span in spans)
            flags = max(int(span.get("flags", 0)) for span in spans)
            sizes.append(size)
            lines.append({
                "text": line_text,
                "size": size,
                "bold": bool(flags & 16),
                "y": min(span.get("bbox", [0, 0, 0, 0])[1] for span in spans),
            })

    median_size = sorted(sizes)[len(sizes) // 2] if sizes else 0
    return {
        "page_index": page_index,
        "text": text,
        "lines": lines,
        "median_size": median_size,
    }


def _detect_pdf_topic_boundaries(page_infos: list[dict]):
    boundaries = []
    seen_titles = set()

    for page in page_infos:
        candidate = _best_heading_for_page(page)
        if not candidate:
            continue
        normalized = _heading_key(candidate)
        if normalized in seen_titles:
            continue
        seen_titles.add(normalized)
        boundaries.append({
            "page_index": page["page_index"],
            "title": candidate,
        })

    if boundaries and boundaries[0]["page_index"] <= 1 and len(boundaries) > 1:
        boundaries = boundaries[1:]

    return boundaries


def _best_heading_for_page(page: dict):
    candidates = []
    for line in page["lines"][:18]:
        text = line["text"]
        if not _looks_like_lesson_heading(text):
            continue
        score = 0
        if text.upper() == text and any(char.isalpha() for char in text):
            score += 3
        if line["bold"]:
            score += 2
        if line["size"] >= page["median_size"] + 1.5:
            score += 2
        if line["y"] < 220:
            score += 1
        if ":" in text:
            score += 1
        candidates.append((score, line["y"], text))

    if not candidates:
        return None

    candidates.sort(key=lambda item: (-item[0], item[1]))
    return candidates[0][2]


def _looks_like_lesson_heading(text: str):
    clean = text.strip(" :-")
    words = clean.split()
    if len(clean) < 8 or len(clean) > 120:
        return False
    if len(words) < 2 or len(words) > 14:
        return False
    if clean.endswith("."):
        return False
    lowered = clean.lower()
    ignored_terms = {
        "republic act", "copyright", "government of the philippines", "department of education",
        "development team", "editors", "author", "writer", "reviewer", "illustrator",
        "management team", "answer key", "references", "table of contents",
    }
    if any(term in lowered for term in ignored_terms):
        return False
    if re.search(r"^[\d\s\-–—]+$", clean):
        return False
    uppercase_ratio = sum(1 for char in clean if char.isupper()) / max(sum(1 for char in clean if char.isalpha()), 1)
    return uppercase_ratio > 0.45 or ":" in clean


def _fallback_pdf_boundaries(page_infos: list[dict]):
    content_pages = [page for page in page_infos if not _is_front_matter_page(page["text"])]
    if not content_pages:
        content_pages = page_infos
    chunk_size = 4
    boundaries = []
    for index in range(0, len(content_pages), chunk_size):
        page = content_pages[index]
        title = _first_clean_line(page["text"]) or f"Topic {len(boundaries) + 1}"
        boundaries.append({"page_index": page["page_index"], "title": title[:120]})
    return boundaries


def _is_front_matter_page(text: str):
    lowered = text.lower()
    front_terms = ["republic act", "copyright", "development team", "department of education", "editors"]
    return any(term in lowered for term in front_terms)


def _first_clean_line(text: str):
    for line in text.splitlines():
        clean = _normalize_pdf_line(line)
        if _looks_like_lesson_heading(clean):
            return clean
    return None


def _render_pdf_topic_pages(document, start: int, end: int):
    PDF_PAGE_DIR.mkdir(parents=True, exist_ok=True)
    urls = []
    for page_index in range(start, end):
        page = document.load_page(page_index)
        image_name = f"{uuid4().hex}_p{page_index + 1}.png"
        destination = PDF_PAGE_DIR / image_name
        pixmap = page.get_pixmap(matrix=fitz_matrix(), alpha=False)
        pixmap.save(str(destination))
        urls.append(f"/static/pdf_pages/{image_name}")
    return urls


def fitz_matrix():
    try:
        import fitz
        return fitz.Matrix(1.6, 1.6)
    except ImportError:
        return None


def _summarize_topic(title: str, content: str):
    sentences = re.split(r"(?<=[.!?])\s+", " ".join(content.split()))
    useful = [sentence for sentence in sentences if len(sentence.split()) >= 6 and title.lower() not in sentence.lower()]
    selected = useful[:2] or sentences[:2]
    summary = " ".join(selected).strip()
    if summary:
        return summary[:420]
    return f"This lesson covers {title} and the key ideas students need to understand."


def _normalize_pdf_line(text: str):
    return re.sub(r"\s+", " ", text.replace("\u00a0", " ")).strip()


def _heading_key(text: str):
    return re.sub(r"[^a-z0-9]+", "", text.lower())


def _extract_material_text(path: Path, title: str, description: str):
    extension = path.suffix.lower()
    try:
        if extension == ".pdf":
            text = _extract_pdf_text(path)
        elif extension == ".docx":
            text = _extract_docx_text(path)
        elif extension == ".pptx":
            text = _extract_pptx_text(path)
        else:
            text = ""
    except Exception:
        text = ""

    if text.strip():
        return text
    return _build_fallback_text(title, description, path.name)


def _extract_pdf_text(path: Path):
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return ""

    reader = PdfReader(str(path))
    return "\n\n".join((page.extract_text() or "").strip() for page in reader.pages)


def _extract_docx_text(path: Path):
    try:
        from docx import Document
    except ImportError:
        return ""

    document = Document(str(path))
    lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n\n".join(lines)


def _extract_pptx_text(path: Path):
    try:
        from pptx import Presentation
    except ImportError:
        return ""

    presentation = Presentation(str(path))
    lines = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            lines.append(f"Slide {slide_number}\n" + "\n".join(slide_lines))
    return "\n\n".join(lines)


def _extract_material_images(path: Path):
    if path.suffix.lower() not in {".docx", ".pptx"}:
        return []

    media_prefix = "word/media/" if path.suffix.lower() == ".docx" else "ppt/media/"
    urls = []
    MATERIAL_IMAGE_DIR.mkdir(parents=True, exist_ok=True)

    try:
        with zipfile.ZipFile(path) as archive:
            for member in archive.namelist():
                if not member.startswith(media_prefix):
                    continue
                extension = Path(member).suffix.lower()
                if extension not in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
                    continue
                image_name = f"{uuid4().hex}{extension}"
                destination = MATERIAL_IMAGE_DIR / image_name
                with archive.open(member) as source, destination.open("wb") as target:
                    shutil.copyfileobj(source, target)
                urls.append(f"/static/material_images/{image_name}")
    except Exception:
        return []

    return urls
